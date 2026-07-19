from pathlib import Path
import csv

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent
SLIDES_DIR = ROOT / "slides"
OUTPUT = ROOT / "Common Circuit Blocks - Slides.pptx"
VIDEO_URL = "https://www.youtube.com/watch?v=DboxH6EYa40"

with (SLIDES_DIR / "index.csv").open(encoding="utf-8", newline="") as handle:
    rows = list(csv.DictReader(handle))

prs = Presentation()
prs.slide_width = Inches(13.333333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

for row in rows:
    slide = prs.slides.add_slide(blank)
    image = SLIDES_DIR / row["file"]
    # The downloaded source is 1920x1080, matching the 16:9 presentation.
    slide.shapes.add_picture(
        str(image), 0, 0, width=prs.slide_width, height=prs.slide_height
    )

    # Add a small readable timestamp linked to the corresponding video position.
    seconds = int(row["timestamp_seconds"])
    box = slide.shapes.add_textbox(
        Inches(11.72), Inches(7.12), Inches(1.48), Inches(0.27)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    box.fill.transparency = 12
    box.line.color.rgb = RGBColor(100, 100, 100)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = row["timestamp"]
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(25, 25, 25)
    run.hyperlink.address = f"{VIDEO_URL}&t={seconds}s"
    paragraph.alignment = PP_ALIGN.CENTER

# PowerPoint starts with one slide only when explicitly added, so no cleanup is needed.
prs.core_properties.title = "Common Circuit Blocks"
prs.core_properties.subject = "Slides reconstructed from Yifan Sun's lecture video"
prs.core_properties.author = "Yifan Sun / slide capture compilation"
prs.core_properties.comments = (
    "Slide images captured from the downloaded lecture video. "
    "Timestamp labels link to the corresponding YouTube positions."
)
prs.save(OUTPUT)
print(f"Created {OUTPUT} with {len(rows)} slides")
